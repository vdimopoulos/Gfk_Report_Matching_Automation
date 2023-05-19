from pptx import Presentation
from collections import defaultdict
import re

class ReportListCreatorTest:
    def extract_codes_from_pptx(source_path):
        
        prs = Presentation(source_path)
        sheets_dict = defaultdict(set)

        RG_pattern = r"\bRG\b\s*(\S+)"
        ID_pattern = r"\bID\b\s*(\S+)"
        RG_found = False
        ID_found = False

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    text_frame = shape.text_frame
                    for paragraph in text_frame.paragraphs:
                        for run in paragraph.runs:
                            text = run.text.strip()
                            RG_match = re.search(RG_pattern, text)
                            ID_match = re.search(ID_pattern, text)
                            if RG_match:
                                RG_text = RG_match.group(1)
                                RG_found = True
                            if ID_match:
                                ID_text = ID_match.group(1)
                                ID_found = True
                    if RG_found and ID_found:
                        sheets_dict[RG_text].add(ID_text)
        
        output_string = ""

        for key, values in sheets_dict.items():
            output_string += f"Report Group <{key}>:\n"

            for i, value in enumerate(values, 1):
                output_string += f"Sheet <{value}>"

                # Add a comma and line break every 5 values
                if i % 5 == 0:
                    output_string += ",\n"
                else:
                    output_string += ", "
            
            output_string += "\n\n"

            # Remove the trailing line breaks
        output_string = output_string.rstrip("\n\n")

        return(output_string)



