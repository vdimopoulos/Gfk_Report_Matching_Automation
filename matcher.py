from pptx import Presentation
import re
import win32com.client as win32



class Matcher():
    
    def merge_ppts(template_path, deck_path, destination_path):

        template_path = template_path.replace('/','\\')
        deck_path = deck_path.replace('/','\\')
        destination_path = destination_path.replace('/','\\')

        def find_id_in_slide(slide):
            # Iterate through shapes in the slide
            for shape in slide.Shapes:
                if shape.HasTextFrame:
                    text_frame = shape.TextFrame
                    if text_frame.HasText:
                        # Get the text content from the shape
                        text = text_frame.TextRange.Text
                        ID_pattern = r"\bID\b\s*(\S+)"
                        ID_match = re.search(ID_pattern, text)
                        if ID_match:
                            # Extract the ID value from the text
                            return ID_match.group(1)
            return None
        

        def copy_slide(slide, target_presentation):
            new_slide_layout = slide.slide_layout
            new_slide = target_presentation.slides.add_slide(new_slide_layout)

            # Copy the source slide's content and formatting to the new slide
            for shape in slide.shapes:
                new_shape = new_slide.shapes.add_shape(shape.auto_shape_type, shape.left, shape.top, shape.width, shape.height)
                new_shape.text_frame.text = shape.text

                # Copy additional formatting properties if needed
                new_shape.fill = shape.fill
                new_shape.line = shape.line

            return new_slide
        


        def merge_presentations(pptx1_path, pptx2_path, output_path):
            ppt_app = win32.Dispatch("PowerPoint.Application")
            
            # Open the first presentation
            pptx1 = ppt_app.Presentations.Open(pptx1_path)
            pptx1_slides_count = pptx1.Slides.Count
            
            # Open the second presentation
            pptx2 = ppt_app.Presentations.Open(pptx2_path)
            
            # Create a new presentation
            pptx3 = ppt_app.Presentations.Add()
            
            # Iterate through slides in the first presentation
            for slide_index in range(1, pptx1_slides_count + 1):
                slide = pptx1.Slides(slide_index)
                
                # Find the ID within the slide
                slide_id = find_id_in_slide(slide)
                
                if slide_id:
                    # Search for the corresponding slide in pptx2
                    for pptx2_slide_index in range(1, pptx2.Slides.Count + 1):
                        pptx2_slide = pptx2.Slides(pptx2_slide_index)
                        pptx2_slide_id = find_id_in_slide(pptx2_slide)

                        if pptx2_slide_id == slide_id:
                            # Copy the slide from pptx2 to pptx3
                            copy_slide(pptx2_slide, pptx3)
                            break
                    else:
                        # If no matching slide is found in pptx2, copy the slide from pptx1 to pptx3
                        copy_slide(slide, pptx3)
                else:
                    # Copy the slide from pptx1 to pptx3
                    copy_slide(slide, pptx3)

            # Save the merged presentation
            pptx3.SaveAs(output_path)
            
            # Close all presentations
            pptx1.Close()
            pptx2.Close()
            pptx3.Close()
            
            # Quit PowerPoint
            ppt_app.Quit()
        
        merge_presentations(template_path, deck_path, destination_path +'\\matched.pptx')
    